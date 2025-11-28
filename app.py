from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.staticfiles import StaticFiles
from fastapi.responses import HTMLResponse, JSONResponse, StreamingResponse
from pydantic import BaseModel
from typing import List, Optional, AsyncGenerator
import os
import sys
import asyncio
from pathlib import Path
import re
from datetime import datetime
import json
import time

# Import AI features
try:
    from ai_features import ai_features
    AI_FEATURES_AVAILABLE = True
except ImportError:
    AI_FEATURES_AVAILABLE = False
    ai_features = None

# File type handlers
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

app = FastAPI(title="Anvesh - Advanced File Search")

# Helper function to get resource path (works for both script and executable)
def resource_path(relative_path):
    """Get absolute path to resource, works for dev and PyInstaller"""
    if getattr(sys, 'frozen', False):
        # Running as compiled executable - PyInstaller creates a temp folder and stores path in _MEIPASS
        base_path = sys._MEIPASS
    else:
        # Running as script
        base_path = os.path.dirname(os.path.abspath(__file__))
    return os.path.join(base_path, relative_path)

# Mount static files
static_path = resource_path("static")
if os.path.exists(static_path):
    app.mount("/static", StaticFiles(directory=static_path), name="static")

class SearchRequest(BaseModel):
    query: str
    folders: List[str]
    exact_match: bool = False
    case_sensitive: bool = False
    search_filenames: bool = False

class SearchResult(BaseModel):
    file_path: str
    line_number: Optional[int]
    content: str
    occurrences: int

class FileResult(BaseModel):
    file_path: str
    total_occurrences: int
    matches: List[SearchResult]

def search_txt(file_path: str, query: str, exact_match: bool, case_sensitive: bool) -> List[SearchResult]:
    """Search in plain text files"""
    results = []
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            lines = f.readlines()
            for line_num, line in enumerate(lines, 1):
                if case_sensitive:
                    search_line = line
                    search_query = query
                else:
                    search_line = line.lower()
                    search_query = query.lower()
                
                if exact_match:
                    if search_query in search_line:
                        results.append(SearchResult(
                            file_path=file_path,
                            line_number=line_num,
                            content=line.strip(),
                            occurrences=search_line.count(search_query)
                        ))
                else:
                    if search_query in search_line:
                        results.append(SearchResult(
                            file_path=file_path,
                            line_number=line_num,
                            content=line.strip()[:200],  # Limit preview
                            occurrences=search_line.count(search_query)
                        ))
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
    return results

def search_docx(file_path: str, query: str, exact_match: bool, case_sensitive: bool) -> List[SearchResult]:
    """Search in Word documents"""
    if not DOCX_AVAILABLE:
        return []
    
    results = []
    try:
        doc = Document(file_path)
        line_num = 0
        for paragraph in doc.paragraphs:
            line_num += 1
            text = paragraph.text
            if case_sensitive:
                search_text = text
                search_query = query
            else:
                search_text = text.lower()
                search_query = query.lower()
            
            if search_query in search_text:
                count = search_text.count(search_query)
                results.append(SearchResult(
                    file_path=file_path,
                    line_number=line_num,
                    content=text.strip()[:200],
                    occurrences=count
                ))
    except Exception as e:
        print(f"Error reading DOCX {file_path}: {e}")
    return results

def search_xlsx(file_path: str, query: str, exact_match: bool, case_sensitive: bool) -> List[SearchResult]:
    """Search in Excel files"""
    if not XLSX_AVAILABLE:
        return []
    
    results = []
    try:
        wb = load_workbook(file_path, data_only=True)
        line_num = 0
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        line_num += 1
                        text = str(cell.value)
                        if case_sensitive:
                            search_text = text
                            search_query = query
                        else:
                            search_text = text.lower()
                            search_query = query.lower()
                        
                        if search_query in search_text:
                            count = search_text.count(search_query)
                            results.append(SearchResult(
                                file_path=f"{file_path} (Sheet: {sheet_name})",
                                line_number=line_num,
                                content=text.strip()[:200],
                                occurrences=count
                            ))
    except Exception as e:
        print(f"Error reading XLSX {file_path}: {e}")
    return results

def search_pptx(file_path: str, query: str, exact_match: bool, case_sensitive: bool) -> List[SearchResult]:
    """Search in PowerPoint files"""
    if not PPTX_AVAILABLE:
        return []
    
    results = []
    try:
        prs = Presentation(file_path)
        slide_num = 0
        for slide in prs.slides:
            slide_num += 1
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text
                    if case_sensitive:
                        search_text = text
                        search_query = query
                    else:
                        search_text = text.lower()
                        search_query = query.lower()
                    
                    if search_query in search_text:
                        count = search_text.count(search_query)
                        results.append(SearchResult(
                            file_path=f"{file_path} (Slide: {slide_num})",
                            line_number=slide_num,
                            content=text.strip()[:200],
                            occurrences=count
                        ))
    except Exception as e:
        print(f"Error reading PPTX {file_path}: {e}")
    return results

def search_pdf(file_path: str, query: str, exact_match: bool, case_sensitive: bool) -> List[SearchResult]:
    """Search in PDF files"""
    if not PDF_AVAILABLE:
        return []
    
    results = []
    try:
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            page_num = 0
            for page in pdf_reader.pages:
                page_num += 1
                text = page.extract_text()
                if case_sensitive:
                    search_text = text
                    search_query = query
                else:
                    search_text = text.lower()
                    search_query = query.lower()
                
                if search_query in search_text:
                    lines = text.split('\n')
                    for line_idx, line in enumerate(lines, 1):
                        if case_sensitive:
                            search_line = line
                        else:
                            search_line = line.lower()
                        
                        if search_query in search_line:
                            count = search_line.count(search_query)
                            results.append(SearchResult(
                                file_path=f"{file_path} (Page: {page_num})",
                                line_number=line_idx,
                                content=line.strip()[:200],
                                occurrences=count
                            ))
    except Exception as e:
        print(f"Error reading PDF {file_path}: {e}")
    return results

def get_supported_files(folder_path: str) -> List[str]:
    """Get all supported files from a folder recursively"""
    supported_extensions = {'.txt', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.pdf'}
    files = []
    
    try:
        for root, dirs, filenames in os.walk(folder_path):
            for filename in filenames:
                ext = os.path.splitext(filename)[1].lower()
                if ext in supported_extensions:
                    file_path = os.path.join(root, filename)
                    files.append(file_path)
    except Exception as e:
        print(f"Error walking directory {folder_path}: {e}")
    
    return files

def log_search_history(query: str, folders: List[str], results_count: int, exact_match: bool, case_sensitive: bool, search_filenames: bool):
    """Log search history to local file"""
    # History file should be in the same directory as the executable (for persistence)
    if getattr(sys, 'frozen', False):
        # Running as executable - save history next to the .exe file
        base_path = os.path.dirname(sys.executable)
    else:
        # Running as script - save in script directory
        base_path = os.path.dirname(os.path.abspath(__file__))
    
    log_file = os.path.join(base_path, "search_history.json")
    history_entry = {
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "query": query,
        "folders": folders,
        "results_count": results_count,
        "exact_match": exact_match,
        "case_sensitive": case_sensitive,
        "search_filenames": search_filenames
    }
    
    try:
        # Read existing history
        if os.path.exists(log_file):
            with open(log_file, 'r', encoding='utf-8') as f:
                history = json.load(f)
        else:
            history = []
        
        # Append new entry
        history.append(history_entry)
        
        # Keep only last 100 entries
        if len(history) > 100:
            history = history[-100:]
        
        # Write back
        with open(log_file, 'w', encoding='utf-8') as f:
            json.dump(history, f, indent=2, ensure_ascii=False)
    except Exception as e:
        print(f"Error logging search history: {e}")

async def search_files_streaming(search_request: SearchRequest):
    """Search for query in all files and yield results as they're found"""
    all_files = []
    
    # Collect all files from selected folders
    for folder in search_request.folders:
        if os.path.isdir(folder):
            all_files.extend(get_supported_files(folder))
    
    total_files = len(all_files)
    files_processed = 0
    results_count = 0
    
    # Send initial status
    yield f"data: {json.dumps({'type': 'status', 'total_files': total_files, 'files_processed': 0, 'message': f'Found {total_files} files to search...'})}\n\n"
    
    # Search in each file
    for file_path in all_files:
        files_processed += 1
        ext = os.path.splitext(file_path)[1].lower()
        matches = []
        
        # If filename search was requested, check filename match
        if search_request.search_filenames:
            filename = os.path.basename(file_path)
            if search_request.case_sensitive:
                search_filename = filename
                search_query = search_request.query
            else:
                search_filename = filename.lower()
                search_query = search_request.query.lower()
            
            if search_query in search_filename:
                count = search_filename.count(search_query)
                matches.append(SearchResult(
                    file_path=file_path,
                    line_number=None,
                    content=f"Filename match: {filename}",
                    occurrences=count
                ))
        
        # Always search file contents
        if ext == '.txt':
            content_matches = search_txt(file_path, search_request.query, search_request.exact_match, search_request.case_sensitive)
            matches.extend(content_matches)
        elif ext in ['.docx', '.doc']:
            content_matches = search_docx(file_path, search_request.query, search_request.exact_match, search_request.case_sensitive)
            matches.extend(content_matches)
        elif ext in ['.xlsx', '.xls']:
            content_matches = search_xlsx(file_path, search_request.query, search_request.exact_match, search_request.case_sensitive)
            matches.extend(content_matches)
        elif ext in ['.pptx', '.ppt']:
            content_matches = search_pptx(file_path, search_request.query, search_request.exact_match, search_request.case_sensitive)
            matches.extend(content_matches)
        elif ext == '.pdf':
            content_matches = search_pdf(file_path, search_request.query, search_request.exact_match, search_request.case_sensitive)
            matches.extend(content_matches)
        
        # If file has matches, send it immediately
        if matches:
            results_count += 1
            total_occurrences = sum(m.occurrences for m in matches)
            file_result = FileResult(
                file_path=file_path,
                total_occurrences=total_occurrences,
                matches=matches
            )
            yield f"data: {json.dumps({'type': 'result', 'data': file_result.model_dump()})}\n\n"
        
        # Send progress update every 10 files or on last file
        if files_processed % 10 == 0 or files_processed == total_files:
            progress = int((files_processed / total_files) * 100) if total_files > 0 else 0
            yield f"data: {json.dumps({'type': 'progress', 'files_processed': files_processed, 'total_files': total_files, 'progress': progress, 'results_found': results_count})}\n\n"
    
    # Send completion
    yield f"data: {json.dumps({'type': 'complete', 'total_results': results_count})}\n\n"

async def search_files(search_request: SearchRequest) -> List[FileResult]:
    """Search for query in all files across selected folders (non-streaming version for compatibility)"""
    all_files = []
    
    # Collect all files from selected folders
    for folder in search_request.folders:
        if os.path.isdir(folder):
            all_files.extend(get_supported_files(folder))
    
    file_results = []
    
    # Search in each file
    for file_path in all_files:
        ext = os.path.splitext(file_path)[1].lower()
        matches = []
        
        # If filename search was requested, check filename match
        if search_request.search_filenames:
            filename = os.path.basename(file_path)
            if search_request.case_sensitive:
                search_filename = filename
                search_query = search_request.query
            else:
                search_filename = filename.lower()
                search_query = search_request.query.lower()
            
            if search_query in search_filename:
                count = search_filename.count(search_query)
                matches.append(SearchResult(
                    file_path=file_path,
                    line_number=None,
                    content=f"Filename match: {filename}",
                    occurrences=count
                ))
        
        # Always search file contents
        if ext == '.txt':
            content_matches = search_txt(file_path, search_request.query, search_request.exact_match, search_request.case_sensitive)
            matches.extend(content_matches)
        elif ext in ['.docx', '.doc']:
            content_matches = search_docx(file_path, search_request.query, search_request.exact_match, search_request.case_sensitive)
            matches.extend(content_matches)
        elif ext in ['.xlsx', '.xls']:
            content_matches = search_xlsx(file_path, search_request.query, search_request.exact_match, search_request.case_sensitive)
            matches.extend(content_matches)
        elif ext in ['.pptx', '.ppt']:
            content_matches = search_pptx(file_path, search_request.query, search_request.exact_match, search_request.case_sensitive)
            matches.extend(content_matches)
        elif ext == '.pdf':
            content_matches = search_pdf(file_path, search_request.query, search_request.exact_match, search_request.case_sensitive)
            matches.extend(content_matches)
        
        if matches:
            total_occurrences = sum(m.occurrences for m in matches)
            file_results.append(FileResult(
                file_path=file_path,
                total_occurrences=total_occurrences,
                matches=matches
            ))
    
    return file_results

@app.get("/", response_class=HTMLResponse)
async def read_root():
    """Serve the main HTML page"""
    template_path = resource_path("templates/index.html")
    if not os.path.exists(template_path):
        return HTMLResponse(content="<h1>Error: templates/index.html not found</h1><p>Path: " + template_path + "</p>", status_code=500)
    with open(template_path, "r", encoding="utf-8") as f:
        return HTMLResponse(content=f.read())

@app.post("/api/search")
async def search(search_request: SearchRequest):
    """Search endpoint - streaming results"""
    async def generate():
        start_time = time.time()
        results_count = 0
        
        async for chunk in search_files_streaming(search_request):
            yield chunk
            # Parse chunk to get results count
            if chunk.startswith("data: "):
                try:
                    data = json.loads(chunk[6:])
                    if data.get('type') == 'result':
                        results_count += 1
                    elif data.get('type') == 'complete':
                        results_count = data.get('total_results', results_count)
                        # Log search history after completion
                        log_search_history(
                            search_request.query,
                            search_request.folders,
                            results_count,
                            search_request.exact_match,
                            search_request.case_sensitive,
                            search_request.search_filenames
                        )
                except:
                    pass
    
    return StreamingResponse(generate(), media_type="text/event-stream")

@app.post("/api/search-sync")
async def search_sync(search_request: SearchRequest):
    """Synchronous search endpoint (for compatibility)"""
    try:
        results = await search_files(search_request)
        results_count = len(results)
        
        # Log search history
        log_search_history(
            search_request.query,
            search_request.folders,
            results_count,
            search_request.exact_match,
            search_request.case_sensitive,
            search_request.search_filenames
        )
        
        return JSONResponse(content={"results": [r.model_dump() for r in results]})
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/history")
async def get_history():
    """Get search history"""
    # History file should be in the same directory as the executable (for persistence)
    if getattr(sys, 'frozen', False):
        # Running as executable - save history next to the .exe file
        base_path = os.path.dirname(sys.executable)
    else:
        # Running as script - save in script directory
        base_path = os.path.dirname(os.path.abspath(__file__))
    
    log_file = os.path.join(base_path, "search_history.json")
    try:
        if os.path.exists(log_file):
            with open(log_file, 'r', encoding='utf-8') as f:
                history = json.load(f)
            # Return in reverse order (newest first)
            return JSONResponse(content={"history": list(reversed(history))})
        else:
            return JSONResponse(content={"history": []})
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/open-file")
async def open_file(request: dict):
    """Open file with default application, optionally at specific line"""
    import subprocess
    import platform
    import shutil
    
    file_path = request.get("path")
    line_number = request.get("line")
    
    if not file_path or not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")
    
    try:
        system = platform.system()
        file_ext = os.path.splitext(file_path)[1].lower()
        
        # Text file extensions that support line numbers
        text_extensions = ('.txt', '.py', '.js', '.html', '.css', '.json', '.md', '.xml', 
                          '.yaml', '.yml', '.ini', '.cfg', '.conf', '.log', '.bat', '.ps1',
                          '.sh', '.sql', '.java', '.cpp', '.c', '.h', '.hpp', '.cs', '.php',
                          '.rb', '.go', '.rs', '.swift', '.kt', '.scala', '.r', '.m', '.pl')
        
        if system == "Windows":
            opened = False
            
            # Try to open with line number support for text files
            if line_number and file_ext in text_extensions:
                # Try VS Code - check in common installation paths and PATH
                vscode_paths = [
                    r'C:\Program Files\Microsoft VS Code\Code.exe',
                    r'C:\Program Files (x86)\Microsoft VS Code\Code.exe',
                    os.path.expanduser(r'~\AppData\Local\Programs\Microsoft VS Code\Code.exe')
                ]
                
                # Check if code is in PATH (but don't use shell=True to avoid error messages)
                code_in_path = shutil.which('code')
                if code_in_path:
                    vscode_paths.insert(0, code_in_path)
                
                # Try VS Code from known paths
                for vscode_path in vscode_paths:
                    if os.path.exists(vscode_path) or (code_in_path and vscode_path == code_in_path):
                        try:
                            # Use CREATE_NO_WINDOW flag to prevent console window
                            subprocess.Popen([vscode_path, '--goto', f'{file_path}:{line_number}'], 
                                            stdout=subprocess.DEVNULL, 
                                            stderr=subprocess.DEVNULL,
                                            creationflags=subprocess.CREATE_NO_WINDOW)
                            opened = True
                            break
                        except:
                            pass
                
                # Try Notepad++ if VS Code failed
                if not opened:
                    npp_paths = [
                        r'C:\Program Files\Notepad++\notepad++.exe',
                        r'C:\Program Files (x86)\Notepad++\notepad++.exe',
                        os.path.expanduser(r'~\AppData\Local\Programs\Notepad++\notepad++.exe')
                    ]
                    for npp_path in npp_paths:
                        if os.path.exists(npp_path):
                            try:
                                subprocess.Popen([npp_path, file_path, '-n' + str(line_number)], 
                                                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                                opened = True
                                break
                            except:
                                pass
                
                # Try Sublime Text
                if not opened:
                    sublime_paths = [
                        r'C:\Program Files\Sublime Text\sublime_text.exe',
                        r'C:\Program Files (x86)\Sublime Text\sublime_text.exe'
                    ]
                    for sublime_path in sublime_paths:
                        if os.path.exists(sublime_path):
                            try:
                                subprocess.Popen([sublime_path, f'{file_path}:{line_number}'], 
                                                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                                opened = True
                                break
                            except:
                                pass
            
            # If no line number or couldn't open with editor, use default application
            if not opened:
                try:
                    os.startfile(file_path)
                    opened = True
                except Exception as e:
                    # Last resort: try with subprocess start command (Windows only)
                    try:
                        # Use CREATE_NO_WINDOW to prevent console window from appearing
                        subprocess.Popen(['cmd', '/c', 'start', '', file_path], 
                                        stdout=subprocess.DEVNULL, 
                                        stderr=subprocess.DEVNULL,
                                        creationflags=subprocess.CREATE_NO_WINDOW)
                        opened = True
                    except:
                        raise HTTPException(status_code=500, detail=f"Could not open file: {str(e)}")
            
            if opened:
                return {"status": "ok", "message": "File opened"}
            else:
                raise HTTPException(status_code=500, detail="Could not open file with any available application")
                
        elif system == "Darwin":  # macOS
            if line_number and file_ext in text_extensions:
                # Try VS Code - check in common installation paths
                vscode_paths = [
                    '/Applications/Visual Studio Code.app/Contents/Resources/app/bin/code',
                    os.path.expanduser('~/Applications/Visual Studio Code.app/Contents/Resources/app/bin/code')
                ]
                code_in_path = shutil.which('code')
                if code_in_path:
                    vscode_paths.insert(0, code_in_path)
                
                for vscode_path in vscode_paths:
                    if os.path.exists(vscode_path) or (code_in_path and vscode_path == code_in_path):
                        try:
                            subprocess.Popen([vscode_path, '--goto', f'{file_path}:{line_number}'],
                                            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                            return {"status": "ok", "message": "File opened"}
                        except:
                            pass
                        break
                # Fallback to TextEdit
                subprocess.Popen(['open', '-a', 'TextEdit', file_path],
                                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            else:
                subprocess.Popen(['open', file_path],
                                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            return {"status": "ok", "message": "File opened"}
            
        else:  # Linux
            if line_number and file_ext in text_extensions:
                # Try VS Code - check in common installation paths
                vscode_paths = [
                    '/usr/bin/code',
                    '/usr/local/bin/code',
                    os.path.expanduser('~/.local/share/code/code')
                ]
                code_in_path = shutil.which('code')
                if code_in_path:
                    vscode_paths.insert(0, code_in_path)
                
                for vscode_path in vscode_paths:
                    if os.path.exists(vscode_path) or (code_in_path and vscode_path == code_in_path):
                        try:
                            subprocess.Popen([vscode_path, '--goto', f'{file_path}:{line_number}'],
                                            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                            return {"status": "ok", "message": "File opened"}
                        except:
                            pass
                        break
            subprocess.Popen(['xdg-open', file_path],
                            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            return {"status": "ok", "message": "File opened"}
            
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error opening file: {str(e)}")

@app.get("/api/health")
async def health():
    """Health check endpoint"""
    health_data = {
        "status": "ok",
        "docx": DOCX_AVAILABLE,
        "xlsx": XLSX_AVAILABLE,
        "pptx": PPTX_AVAILABLE,
        "pdf": PDF_AVAILABLE
    }
    
    if AI_FEATURES_AVAILABLE and ai_features:
        health_data["ai_features"] = ai_features.get_capabilities()
    else:
        health_data["ai_features"] = {"available": False}
    
    return health_data

# ==================== AI Features Routes ====================

@app.get("/ai", response_class=HTMLResponse)
async def ai_features_page():
    """Serve the AI features page"""
    template_path = resource_path("templates/ai_features.html")
    if not os.path.exists(template_path):
        return HTMLResponse(content="<h1>AI Features page not found</h1>", status_code=500)
    with open(template_path, "r", encoding="utf-8") as f:
        return HTMLResponse(content=f.read())

@app.get("/api/ai/capabilities")
async def get_ai_capabilities():
    """Get available AI capabilities"""
    if not AI_FEATURES_AVAILABLE:
        return JSONResponse(content={"available": False, "error": "AI features not available"})
    
    return JSONResponse(content=ai_features.get_capabilities())

class FilePathRequest(BaseModel):
    file_path: str
    frame_interval: Optional[int] = 30
    confidence: Optional[float] = 0.25
    threshold: Optional[float] = 0.6

class FaceMatchRequest(BaseModel):
    image1_path: str
    image2_path: Optional[str] = None
    reference_image_path: Optional[str] = None
    folder_path: Optional[str] = None
    threshold: float = 0.6

@app.post("/api/ai/ocr/image")
async def ocr_image(request: FilePathRequest):
    """Extract text from image using OCR"""
    if not AI_FEATURES_AVAILABLE:
        raise HTTPException(status_code=503, detail="AI features not available")
    
    if not os.path.exists(request.file_path):
        raise HTTPException(status_code=404, detail="File not found")
    
    result = ai_features.extract_text_from_image(request.file_path)
    return JSONResponse(content=result)

@app.post("/api/ai/ocr/video")
async def ocr_video(request: FilePathRequest):
    """Extract text from video using OCR"""
    if not AI_FEATURES_AVAILABLE:
        raise HTTPException(status_code=503, detail="AI features not available")
    
    if not os.path.exists(request.file_path):
        raise HTTPException(status_code=404, detail="File not found")
    
    result = ai_features.extract_text_from_video(request.file_path, request.frame_interval or 30)
    return JSONResponse(content=result)

@app.post("/api/ai/detect/objects")
async def detect_objects(request: FilePathRequest):
    """Detect objects in image"""
    if not AI_FEATURES_AVAILABLE:
        raise HTTPException(status_code=503, detail="AI features not available")
    
    if not os.path.exists(request.file_path):
        raise HTTPException(status_code=404, detail="File not found")
    
    result = ai_features.detect_objects(request.file_path, request.confidence or 0.25)
    return JSONResponse(content=result)

@app.post("/api/ai/detect/faces")
async def detect_faces(request: FilePathRequest):
    """Detect faces in image"""
    if not AI_FEATURES_AVAILABLE:
        raise HTTPException(status_code=503, detail="AI features not available")
    
    if not os.path.exists(request.file_path):
        raise HTTPException(status_code=404, detail="File not found")
    
    result = ai_features.detect_faces(request.file_path)
    return JSONResponse(content=result)

@app.post("/api/ai/detect/faces/video")
async def detect_faces_video(request: FilePathRequest):
    """Detect faces in video"""
    if not AI_FEATURES_AVAILABLE:
        raise HTTPException(status_code=503, detail="AI features not available")
    
    if not os.path.exists(request.file_path):
        raise HTTPException(status_code=404, detail="File not found")
    
    result = ai_features.detect_faces_in_video(request.file_path, request.frame_interval or 30)
    return JSONResponse(content=result)

@app.post("/api/ai/match/faces")
async def match_faces(request: FaceMatchRequest):
    """Match faces between two images"""
    if not AI_FEATURES_AVAILABLE:
        raise HTTPException(status_code=503, detail="AI features not available")
    
    if not request.image2_path:
        raise HTTPException(status_code=400, detail="image2_path is required")
    
    if not os.path.exists(request.image1_path) or not os.path.exists(request.image2_path):
        raise HTTPException(status_code=404, detail="One or both images not found")
    
    result = ai_features.match_faces(request.image1_path, request.image2_path, request.threshold)
    return JSONResponse(content=result)

@app.post("/api/ai/match/faces/folder")
async def find_matching_faces(request: FaceMatchRequest):
    """Find matching faces in folder"""
    if not AI_FEATURES_AVAILABLE:
        raise HTTPException(status_code=503, detail="AI features not available")
    
    if not request.reference_image_path or not request.folder_path:
        raise HTTPException(status_code=400, detail="reference_image_path and folder_path are required")
    
    if not os.path.exists(request.reference_image_path):
        raise HTTPException(status_code=404, detail="Reference image not found")
    
    if not os.path.isdir(request.folder_path):
        raise HTTPException(status_code=404, detail="Folder not found")
    
    result = ai_features.find_matching_faces_in_folder(request.reference_image_path, request.folder_path, request.threshold)
    return JSONResponse(content=result)

if __name__ == "__main__":
    import uvicorn
    import webbrowser
    import threading
    import time
    
    # Detect if running as PyInstaller executable
    if getattr(sys, 'frozen', False):
        # Running as compiled executable
        # Change to executable directory for data files (history, etc.)
        application_path = os.path.dirname(sys.executable)
        os.chdir(application_path)
        print("=" * 60)
        print("  अन्वेष (Anvesh) - Advanced File Search")
        print("=" * 60)
        print("\nStarting server...")
        print("Please wait, this may take a moment...\n")
    else:
        # Running as script
        print("=" * 60)
        print("  अन्वेष (Anvesh) - Advanced File Search")
        print("=" * 60)
        print("\nStarting server on http://127.0.0.1:8000")
        print("Press Ctrl+C to stop the server\n")
    
    # Function to open browser after delay
    def open_browser():
        time.sleep(2)  # Wait for server to start
        try:
            url = "http://127.0.0.1:8000"
            webbrowser.open(url)
            if getattr(sys, 'frozen', False):
                print(f"✓ Browser opened at {url}")
                print("\n" + "=" * 60)
                print("Server is running!")
                print("=" * 60)
                print("\nKeep this window open while using Anvesh.")
                print("Close this window to stop the server.\n")
        except Exception as e:
            print(f"Could not open browser automatically: {e}")
            print(f"Please open http://127.0.0.1:8000 in your browser manually")
    
    # Start browser opener in background
    browser_thread = threading.Thread(target=open_browser, daemon=True)
    browser_thread.start()
    
    try:
        uvicorn.run(app, host="127.0.0.1", port=8000, log_level="info")
    except KeyboardInterrupt:
        print("\n\nServer stopped by user.")
    except Exception as e:
        print(f"\n\nERROR: Failed to start server: {e}")
        print("\nPress any key to exit...")
        if getattr(sys, 'frozen', False):
            input()

